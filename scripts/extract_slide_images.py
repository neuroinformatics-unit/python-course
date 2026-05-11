"""
One-off script: extract images and text from .pptx slide decks into public/images/ and scripts/slide-text/.

Usage:
    python scripts/extract_slide_images.py

Outputs:
    public/images/session1/slide01-01.png  (one PNG per picture shape per slide)
    scripts/slide-text/session1-slide01.md (title + text frames + notes per slide)
    scripts/slide-text/index.json          (lookup: session/slide -> {imagePaths, textPath, title})
"""

import json
import os
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).parent.parent
MATERIALS = ROOT / "materials"
PUBLIC_IMAGES = ROOT / "public" / "images"
SLIDE_TEXT = Path(__file__).parent / "slide-text"

DECKS = [
    ("session1", "Python Club Session 1.pptx"),
    ("session2", "Python Club Session 2.pptx"),
    ("session3", "Python club session 3.pptx"),
    ("session4", "Python club session 4.pptx"),
    ("osw", "software_skills_intro_OSW_2025.pptx"),
]

index = {}


def slugify(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", text.lower()).strip("-")


def extract_text(slide) -> tuple[str, list[str]]:
    title = ""
    paragraphs: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.idx == 0:
            title = text
        else:
            paragraphs.extend(line.strip() for line in text.splitlines() if line.strip())
    notes_text = ""
    if slide.has_notes_slide:
        notes_text = slide.notes_slide.notes_text_frame.text.strip()
    return title, paragraphs, notes_text


def extract_images(slide, out_dir: Path, slide_label: str) -> list[str]:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved: list[str] = []
    img_idx = 1
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            img = shape.image
            ext = img.ext.lower()
            if ext in ("jpeg", "jpg"):
                ext = "jpg"
            filename = f"{slide_label}-{img_idx:02d}.{ext}"
            path = out_dir / filename
            path.write_bytes(img.blob)
            saved.append(str(path.relative_to(ROOT / "public")))
            img_idx += 1
    return saved


for session_key, filename in DECKS:
    pptx_path = MATERIALS / filename
    if not pptx_path.exists():
        print(f"  SKIP (not found): {filename}")
        continue

    print(f"Processing {filename} -> {session_key}")
    prs = Presentation(str(pptx_path))
    img_out_dir = PUBLIC_IMAGES / session_key
    SLIDE_TEXT.mkdir(parents=True, exist_ok=True)

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_label = f"slide{slide_num:02d}"
        title, paragraphs, notes = extract_text(slide)
        image_paths = extract_images(slide, img_out_dir, slide_label)

        # Write markdown text dump
        md_path = SLIDE_TEXT / f"{session_key}-{slide_label}.md"
        with md_path.open("w") as f:
            f.write(f"# {title or '(no title)'}\n\n")
            for p in paragraphs:
                f.write(f"- {p}\n")
            if notes:
                f.write(f"\n## Notes\n\n{notes}\n")

        key = f"{session_key}/{slide_label}"
        index[key] = {
            "title": title,
            "textPath": f"scripts/slide-text/{session_key}-{slide_label}.md",
            "imagePaths": [f"/images/{p}" for p in image_paths],
        }
        if image_paths:
            print(f"  {slide_label}: '{title[:50]}' -> {len(image_paths)} image(s)")

(SLIDE_TEXT / "index.json").write_text(json.dumps(index, indent=2))
print(f"\nDone. {sum(len(v['imagePaths']) for v in index.values())} images extracted.")
print(f"Index written to scripts/slide-text/index.json")
